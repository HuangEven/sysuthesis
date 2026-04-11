from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
FIGURES_DIR = ROOT / "figures"
PPTX_PATH = FIGURES_DIR / "fig1_2_scale_bottleneck_migration.pptx"
PNG_PATH = FIGURES_DIR / "fig1_2_scale_bottleneck_migration.png"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
PNG_W = 1600
PNG_H = 900

FONT_FAMILY = "PingFang SC"
FONT_FILE = "/System/Library/Fonts/PingFang.ttc"


COLORS = {
    "lane_fill": "F5F5F5",
    "lane_line": "9A9A9A",
    "blue_fill": "DAE8FC",
    "blue_line": "6C8EBF",
    "yellow_fill": "FFF2CC",
    "yellow_line": "D6B656",
    "red_fill": "F8CECC",
    "red_line": "B85450",
    "green_fill": "D5E8D4",
    "green_line": "82B366",
    "gray_fill": "E5E5E5",
    "gray_line": "A0A0A0",
    "arrow": "6E6E6E",
    "text": "111111",
}


LANES = [
    {
        "title": "小规模数据",
        "x": 0.16,
        "boxes": [
            ("主要特征\n样本量较小\n单机内存可容纳", "blue"),
            ("典型处理方式\n单机 CPU 检测\n规则与统计方法为主", "yellow"),
            ("主要瓶颈\nCPU 计算开销\n内存访问与算法复杂度", "red"),
            ("研究关注点\n算法效率优化\n特征工程与规则筛查", "green"),
        ],
    },
    {
        "title": "中等规模数据",
        "x": 3.40,
        "boxes": [
            ("主要特征\n样本量继续增加\n单机内存接近上限", "blue"),
            ("典型处理方式\n单机 GPU 加速\n批处理与向量化计算", "yellow"),
            ("主要瓶颈\n显存容量受限\nCPU↔GPU 数据搬运", "red"),
            ("研究关注点\n算子批处理、算子融合\n减少主机与设备往返", "green"),
        ],
    },
    {
        "title": "大规模数据",
        "x": 6.64,
        "boxes": [
            ("主要特征\n数据超出显存容量\n索引与原始数据难以内存常驻", "blue"),
            ("典型处理方式\n分层索引 / 候选过滤\n外存辅助的加速检索", "yellow"),
            ("主要瓶颈\nI/O 带宽与访问粒度\n计算与存储阶段串行化", "red"),
            ("研究关注点\n计算—I/O 协同\n异步流水与访问优化", "green"),
        ],
    },
    {
        "title": "超大规模数据",
        "x": 9.88,
        "boxes": [
            ("主要特征\n百亿级 / 持续增长\n查询并发与吞吐要求更高", "blue"),
            ("典型处理方式\n多 GPU / 分布式执行\n跨设备并行与结果归并", "yellow"),
            ("主要瓶颈\n设备间通信与归并开销\n负载均衡与 I/O 竞争", "red"),
            ("研究关注点\n多设备协同与扩展性\n并行调度与资源利用率", "green"),
        ],
    },
]


def rgb(hex_value: str) -> RGBColor:
    return RGBColor.from_string(hex_value)


def hex_to_rgb(hex_value: str) -> tuple[int, int, int]:
    hex_value = hex_value.lstrip("#")
    return tuple(int(hex_value[i : i + 2], 16) for i in (0, 2, 4))


def make_pptx() -> None:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    title = slide.shapes.add_textbox(Inches(4.95), Inches(0.08), Inches(3.45), Inches(0.30))
    tf = title.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "数据规模持续增长"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = FONT_FAMILY
    r.font.size = Pt(21)
    r.font.bold = True
    r.font.color.rgb = rgb(COLORS["text"])

    arrow = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
        Inches(0.80),
        Inches(0.42),
        Inches(11.70),
        Inches(0.23),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = rgb(COLORS["gray_fill"])
    arrow.line.color.rgb = rgb(COLORS["gray_line"])
    arrow.line.width = Pt(1.5)

    lane_y = 0.62
    lane_w = 3.05
    lane_h = 5.75
    box_x_offset = 0.24
    box_w = 2.58
    box_h = 0.84
    box_y_positions = [1.18, 2.28, 3.38, 4.48]

    lane_shapes = []
    red_shapes = []

    for lane in LANES:
        lane_shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(lane["x"]),
            Inches(lane_y),
            Inches(lane_w),
            Inches(lane_h),
        )
        lane_shape.fill.solid()
        lane_shape.fill.fore_color.rgb = rgb(COLORS["lane_fill"])
        lane_shape.line.color.rgb = rgb(COLORS["lane_line"])
        lane_shape.line.width = Pt(1.75)
        lane_shapes.append(lane_shape)

        title_box = slide.shapes.add_textbox(
            Inches(lane["x"] + 0.40), Inches(0.75), Inches(2.25), Inches(0.34)
        )
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = lane["title"]
        p.alignment = PP_ALIGN.CENTER
        r = p.runs[0]
        r.font.name = FONT_FAMILY
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = rgb(COLORS["text"])

        for (text, palette), y in zip(lane["boxes"], box_y_positions, strict=True):
            shape = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(lane["x"] + box_x_offset),
                Inches(y),
                Inches(box_w),
                Inches(box_h),
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = rgb(COLORS[f"{palette}_fill"])
            shape.line.color.rgb = rgb(COLORS[f"{palette}_line"])
            shape.line.width = Pt(1.5)
            shape.text_frame.clear()
            shape.text_frame.word_wrap = True
            shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
            shape.text_frame.margin_left = Inches(0.12)
            shape.text_frame.margin_right = Inches(0.12)
            shape.text_frame.margin_top = Inches(0.06)
            shape.text_frame.margin_bottom = Inches(0.05)
            p = shape.text_frame.paragraphs[0]
            p.text = text
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FONT_FAMILY
                run.font.size = Pt(15.5)
                run.font.bold = False
                run.font.color.rgb = rgb(COLORS["text"])
            if palette == "red":
                red_shapes.append(shape)

    for src, dst in zip(red_shapes, red_shapes[1:], strict=False):
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            src.left + src.width,
            src.top + src.height / 2,
            dst.left,
            dst.top + dst.height / 2,
        )
        connector.line.color.rgb = rgb(COLORS["arrow"])
        connector.line.width = Pt(1.6)
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    bottom = slide.shapes.add_textbox(Inches(1.20), Inches(6.64), Inches(10.90), Inches(0.34))
    tf = bottom.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "瓶颈迁移：CPU计算 → 显存与数据搬运 → I/O带宽与串行化 → 多设备通信与归并"
    p.alignment = PP_ALIGN.CENTER
    r = p.runs[0]
    r.font.name = FONT_FAMILY
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = rgb(COLORS["text"])

    prs.save(PPTX_PATH)


def _fit_font(size: int) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT_FILE, size)


def draw_centered_multiline(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    text: str,
    font: ImageFont.FreeTypeFont,
    fill: tuple[int, int, int],
    spacing: int,
) -> None:
    x0, y0, x1, y1 = box
    bbox = draw.multiline_textbbox((0, 0), text, font=font, align="center", spacing=spacing)
    tw = bbox[2] - bbox[0]
    th = bbox[3] - bbox[1]
    tx = x0 + (x1 - x0 - tw) / 2
    ty = y0 + (y1 - y0 - th) / 2 - 2
    draw.multiline_text((tx, ty), text, font=font, fill=fill, align="center", spacing=spacing)


def draw_arrow(draw: ImageDraw.ImageDraw, x: int, y: int, w: int, h: int) -> None:
    head = int(h * 2.0)
    points = [
        (x, y),
        (x + w - head, y),
        (x + w - head, y - h // 3),
        (x + w, y + h // 2),
        (x + w - head, y + h + h // 3),
        (x + w - head, y + h),
        (x, y + h),
    ]
    draw.polygon(points, fill=hex_to_rgb(COLORS["gray_fill"]), outline=hex_to_rgb(COLORS["gray_line"]))


def draw_dashed_arrow(
    draw: ImageDraw.ImageDraw, start: tuple[int, int], end: tuple[int, int], color: tuple[int, int, int]
) -> None:
    x1, y1 = start
    x2, y2 = end
    dash = 12
    gap = 10
    cursor = x1
    while cursor < x2 - 18:
        seg_end = min(cursor + dash, x2 - 18)
        draw.line((cursor, y1, seg_end, y2), fill=color, width=3)
        cursor = seg_end + gap
    arrow = [(x2 - 16, y2 - 8), (x2 - 16, y2 + 8), (x2, y2)]
    draw.polygon(arrow, fill=color)


def make_png() -> None:
    image = Image.new("RGB", (PNG_W, PNG_H), "white")
    draw = ImageDraw.Draw(image)

    title_font = _fit_font(46)
    lane_title_font = _fit_font(38)
    block_font = _fit_font(31)
    footer_font = _fit_font(37)

    draw.text((PNG_W / 2, 44), "数据规模持续增长", font=title_font, fill=hex_to_rgb(COLORS["text"]), anchor="mm")
    draw_arrow(draw, 95, 92, 1410, 28)

    lane_y = 122
    lane_w = 365
    lane_h = 690
    lane_gap = 24
    lane_xs = [20 + i * (lane_w + lane_gap) for i in range(4)]

    box_x_pad = 33
    box_w = lane_w - 2 * box_x_pad
    box_h = 118
    box_ys = [182, 320, 458, 596]

    for lane, x in zip(LANES, lane_xs, strict=True):
        draw.rounded_rectangle(
            (x, lane_y, x + lane_w, lane_y + lane_h),
            radius=36,
            fill=hex_to_rgb(COLORS["lane_fill"]),
            outline=hex_to_rgb(COLORS["lane_line"]),
            width=3,
        )
        draw.text(
            (x + lane_w / 2, 155),
            lane["title"],
            font=lane_title_font,
            fill=hex_to_rgb(COLORS["text"]),
            anchor="mm",
        )

        for (text, palette), y in zip(lane["boxes"], box_ys, strict=True):
            draw.rounded_rectangle(
                (x + box_x_pad, y, x + box_x_pad + box_w, y + box_h),
                radius=22,
                fill=hex_to_rgb(COLORS[f"{palette}_fill"]),
                outline=hex_to_rgb(COLORS[f"{palette}_line"]),
                width=3,
            )
            draw_centered_multiline(
                draw,
                (x + box_x_pad + 18, y + 10, x + box_x_pad + box_w - 18, y + box_h - 10),
                text,
                block_font,
                hex_to_rgb(COLORS["text"]),
                spacing=6,
            )

    red_mid_y = box_ys[2] + box_h // 2
    red_rights = [x + box_x_pad + box_w for x in lane_xs]
    red_lefts = [x + box_x_pad for x in lane_xs]
    for i in range(3):
        draw_dashed_arrow(
            draw,
            (red_rights[i] + 6, red_mid_y),
            (red_lefts[i + 1] - 10, red_mid_y),
            hex_to_rgb(COLORS["arrow"]),
        )

    draw.text(
        (PNG_W / 2, 853),
        "瓶颈迁移：CPU计算 → 显存与数据搬运 → I/O带宽与串行化 → 多设备通信与归并",
        font=footer_font,
        fill=hex_to_rgb(COLORS["text"]),
        anchor="mm",
    )

    image.save(PNG_PATH, dpi=(200, 200))


if __name__ == "__main__":
    make_pptx()
    make_png()
